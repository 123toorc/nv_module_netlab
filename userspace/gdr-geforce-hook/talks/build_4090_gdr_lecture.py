#!/usr/bin/env python3
"""4090 打通 GDR 讲座 PPT —— 面向内核背景不多的基础程序员。

生成:
  python3 userspace/gdr-geforce-hook/talks/build_4090_gdr_lecture.py

输出与本脚本同目录: 4090-gdr-ioctl-hook.pptx
"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "4090-gdr-ioctl-hook.pptx"

# 16:9
SW, SH = Inches(13.333), Inches(7.5)
# 讲者机器上更常见：Windows 微软雅黑 / Consolas。eastAsia 与正文都写这个名字。
FONT_CN = "微软雅黑"
FONT_MONO = "Consolas"

BG = RGBColor(0x0B, 0x12, 0x20)
BG_CARD = RGBColor(0x15, 0x1D, 0x2E)
BG_CARD2 = RGBColor(0x1B, 0x27, 0x3C)
BG_CODE = RGBColor(0x08, 0x0C, 0x16)
AMBER = RGBColor(0xF0, 0xB4, 0x29)
BLUE = RGBColor(0x4C, 0xB3, 0xFF)
TEAL = RGBColor(0x3E, 0xCF, 0xB2)
GREEN = RGBColor(0x3E, 0xCF, 0x8E)
CORAL = RGBColor(0xF0, 0x71, 0x78)
PURPLE = RGBColor(0xC3, 0xA6, 0xFF)
WHITE = RGBColor(0xF4, 0xF7, 0xFC)
MUTED = RGBColor(0x8B, 0x9B, 0xB4)
LINE = RGBColor(0x2A, 0x38, 0x52)
SOFT = RGBColor(0xD5, 0xDE, 0xEC)

def _rpr(run):
    return run._r.get_or_add_rPr()


def style_run(run, size, color, bold=False, mono=False, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    face = FONT_MONO if mono else FONT_CN
    run.font.name = face
    rPr = _rpr(run)
    for tag, name in (("a:latin", face), ("a:ea", FONT_CN), ("a:cs", face)):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", FONT_CN if tag == "a:ea" else face)


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.strip() + "\n"
    style_run(run, 14, RGBColor(0x22, 0x22, 0x22))


def box(slide, l, t, w, h, fill=None, line=None, radius=0.08):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill or BG_CARD
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.25)
    return shp


def rect(slide, l, t, w, h, fill, line=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.15)
    return shp


def arrow(slide, l, t, w, h, fill=AMBER):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def tb(slide, l, t, w, h, lines, size=18, color=WHITE, bold=False, align="left",
       mono=False, anchor=MSO_ANCHOR.TOP, spacing=1.08):
    """lines: str or list of (text, size, color, bold) / str."""
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    try:
        tf._txBody.bodyPr.set("anchor", {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }.get(anchor, "t"))
    except Exception:
        pass
    if isinstance(lines, str):
        lines = [lines]
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map[align]
        p.space_after = Pt(2)
        p.line_spacing = spacing
        if isinstance(item, tuple):
            text, sz, col, bd = item[0], item[1], item[2], item[3]
            it_mono = item[4] if len(item) > 4 else mono
        else:
            text, sz, col, bd, it_mono = item, size, color, bold, mono
        run = p.add_run()
        run.text = text
        style_run(run, sz, col, bold=bd, mono=it_mono)
    return tx


def footer(slide, page, total, chip=""):
    rect(slide, 0, 7.28, 13.333, 0.22, RGBColor(0x08, 0x0D, 0x16))
    tb(slide, 0.5, 7.28, 7.4, 0.22,
       chip or "4090 打通 GDR  ·  挟持 ioctl，而不是改内核",
       11, MUTED, anchor=MSO_ANCHOR.MIDDLE)
    tb(slide, 10.4, 7.28, 2.4, 0.22, f"{page}  /  {total}",
       11, MUTED, align="right", anchor=MSO_ANCHOR.MIDDLE)


def header(slide, kicker, title, page, total):
    set_slide_bg(slide)
    rect(slide, 0, 0, 0.12, 7.5, AMBER)
    tb(slide, 0.48, 0.18, 10.5, 0.28, kicker, 13, AMBER, bold=True)
    tb(slide, 0.48, 0.42, 12.3, 0.48, title, 28, WHITE, bold=True)
    rect(slide, 0.50, 0.96, 1.35, 0.045, AMBER)
    footer(slide, page, total, kicker)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def pill(slide, l, t, w, h, text, fill, fg=BG):
    box(slide, l, t, w, h, fill, radius=0.5)
    tb(slide, l, t, w, h, text, 13, fg, bold=True, align="center",
       anchor=MSO_ANCHOR.MIDDLE)


def card_title(slide, l, t, w, h, title, body, accent=AMBER, title_size=16, body_size=15):
    box(slide, l, t, w, h, BG_CARD, LINE, 0.06)
    rect(slide, l, t, 0.09, h, accent)
    tb(slide, l + 0.22, t + 0.12, w - 0.34, 0.34, title, title_size, accent, bold=True)
    body_lines = body.split("\n") if isinstance(body, str) else body
    tb(slide, l + 0.22, t + 0.46, w - 0.34, h - 0.58, body_lines, body_size, SOFT)


def bullets(slide, l, t, w, h, items, size=18, color=SOFT, gap=0.02):
    lines = []
    for it in items:
        if isinstance(it, tuple):
            lines.append(it)
        else:
            lines.append((f"•  {it}", size, color, False))
    tb(slide, l, t, w, h, lines, size, color, spacing=1.18)


def code_card(slide, l, t, w, h, lines, title=None):
    box(slide, l, t, w, h, BG_CODE, RGBColor(0x2C, 0x3E, 0x5A), 0.05)
    y = t + 0.10
    if title:
        tb(slide, l + 0.16, y, w - 0.3, 0.28, title, 12, AMBER, bold=True, mono=True)
        y += 0.30
    tb(slide, l + 0.16, y, w - 0.32, h - (y - t) - 0.10,
       [(ln, 13, SOFT if not ln.strip().startswith("#") and "//" not in ln[:8] else MUTED,
         False, True) if not isinstance(ln, tuple) else ln
        for ln in lines],
       13, SOFT, mono=True, spacing=1.12)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    slides = []

    def add():
        s = new_slide(prs)
        slides.append(s)
        return s

    # ---------- 1 title ----------
    def s_title():
        s = add()
        set_slide_bg(s)
        rect(s, 0, 0, 0.16, 7.5, AMBER)
        rect(s, 0, 7.18, 13.333, 0.32, RGBColor(0x08, 0x0D, 0x16))
        tb(s, 0.7, 0.55, 12, 0.32, "内部技术讲座  ·  RTX 4090D  ·  NVIDIA 575.64.05", 15, AMBER, True)
        tb(s, 0.7, 1.05, 12.2, 1.5, "4090 打通 GPUDirect RDMA", 40, WHITE, True)
        tb(s, 0.7, 2.45, 12.0, 0.7, "从 Harry 的 5090 一字节，到我们挟持 ioctl 的 .so", 24, BLUE)
        tb(s, 0.7, 3.25, 11.5, 0.7,
           "听众不需要会写内核。只要写过 C、调过库、见过「打开设备文件」，就能跟上。",
           16, MUTED)
        for i, (lab, col) in enumerate((
            ("不改 NVIDIA 内核", CORAL),
            ("不改 libcuda 文件", PURPLE),
            ("只给测试进程 LD_PRELOAD", GREEN),
        )):
            pill(s, 0.7 + i * 3.55, 4.35, 3.3, 0.46, lab, col, BG)
        tb(s, 0.7, 5.15, 11.5, 0.9,
           "45–60 分钟  ·  源码 userspace/gdr-geforce-hook/gdr_geforce_hook.c\n"
           "对照文章 https://harrychen.xyz/2026/05/20/enable-gpudirect-rdma-on-rtx-5090/",
           14, MUTED)
        tb(s, 0.7, 7.20, 12, 0.26, "讲者备注写在每页备注栏，用「备注视图 / 演讲者视图」讲。", 12, MUTED)
        add_notes(s, """
开场 40 秒：今天不讲「怎么改 NVIDIA 开源内核把 GDR 打开」。我们 4090 上那条路走不通，而且也不该再走。
今天要讲的是：用户态一个 .so，怎样在 CUDA 和内核中间插一脚，把 GeForce 从来不肯写的「第三方 P2P 登记本」补上。
先声明三件事：不改内核、不改 libcuda 文件、不要全局 export LD_PRELOAD。后面每一页都围着这三句转。
听众如果只听懂一个比喻，就是酒店前台的访客登记本。
""")

    # ---------- 2 takeaways ----------
    def s_goals():
        s = add()
        header(s, "00  开场", "今天结束，你要能用自己的话讲清三件事", len(slides), TOTAL[0])
        cards = [
            ("01", AMBER, "GDR 在省哪一次拷贝",
             "网卡直接 DMA 读 GPU 显存，不再经 CPU 内存中转。这是性能故事，也是后面所有报错的背景。"),
            ("02", BLUE, "为什么 4090「硬件能 DMA」仍然 pin 失败",
             "Tesla 的 libcuda 会自己建一本登记册。GeForce 的 libcuda 不建。内核查无此人，返回 0x57 / -22。"),
            ("03", TEAL, "我们的 .so 怎样挟持 ioctl",
             "偷看 CUDA 发给 NVIDIA RM 的「建对象」工单，记下房间号，再自己补写登记本。之后库存 GDRCopy 就能 pin。"),
        ]
        for i, (num, col, title, body) in enumerate(cards):
            x = 0.48 + i * 4.2
            box(s, x, 1.25, 4.0, 4.55, BG_CARD, LINE, 0.07)
            tb(s, x + 0.25, 1.45, 3.5, 0.45, num, 28, col, True)
            tb(s, x + 0.25, 2.05, 3.5, 1.1, title, 20, WHITE, True)
            tb(s, x + 0.25, 3.25, 3.5, 2.2, body, 15, SOFT)
        add_notes(s, """
用 1 分钟把三件事钉死。后面所有细节都是在填这三格。
强调第 2 点：很多人一上来就说「消费卡硬件不支持 GDR」。我们自己也这么以为过。Harry 的 5090 文章和我们 4090 上的日志都说明：硬件 DMA 路径在，缺的是软件政策。
第 3 点预告手法：不是 patch 内核，是挟持用户态的 ioctl。
""")

    # ---------- 3 contract ----------
    def s_contract():
        s = add()
        header(s, "00  开场", "跟听众的约定：今天不需要的知识", len(slides), TOTAL[0])
        left = [
            "不讲页表怎么走、IOMMU 怎么配",
            "不讲 UMMU / SVA / MATT（那是 GDR 打通之后的故事）",
            "不现场反汇编 libcuda",
            "不 presubmit 一份「去改 NVIDIA 内核」的补丁",
        ]
        right = [
            ("ioctl", "用户态给驱动递的一张工单。打开 /dev/xxx 之后，用一个数字说「请做这件事」。"),
            ("handle", "驱动回给你的房间号。每次运行都不同，不要写死。"),
            ("LD_PRELOAD", "让动态链接器先加载我们的 .so，同名函数会被我们的实现「插队」。"),
            ("登记本", "NVIDIA RM 里一本叫 ThirdPartyP2P 的对象。网卡来 pin 时就查它。"),
        ]
        box(s, 0.48, 1.22, 5.7, 5.55, BG_CARD, LINE)
        tb(s, 0.72, 1.40, 5.2, 0.4, "今天故意不讲", 18, CORAL, True)
        bullets(s, 0.72, 1.95, 5.2, 4.4, left, 16)
        box(s, 6.4, 1.22, 6.4, 5.55, BG_CARD, LINE)
        tb(s, 6.64, 1.40, 5.95, 0.4, "只要记住这四个词", 18, GREEN, True)
        y = 1.95
        for name, desc in right:
            tb(s, 6.64, y, 1.6, 0.85, name, 15, AMBER, True, mono=True)
            tb(s, 8.3, y, 4.25, 0.85, desc, 14, SOFT)
            y += 1.1
        add_notes(s, """
这一页是给「我没写过驱动」的人吃定心丸。
把四个词写在白板上也行：ioctl、handle、LD_PRELOAD、登记本。后面反复用。
UMMU 那一页最后只会闪一下，避免把讲座带跑。
""")

    # ---------- 4 bounce vs gdr ----------
    def s_gdr_picture():
        s = add()
        header(s, "01  背景", "先记住一张图：GDR 在省哪一次拷贝", len(slides), TOTAL[0])
        # left: bounce
        box(s, 0.48, 1.22, 6.05, 5.55, BG_CARD, LINE)
        tb(s, 0.7, 1.38, 5.6, 0.36, "普通路径：CPU 中转", 18, CORAL, True)
        nodes = [
            (0.75, 2.00, "GPU 显存", BLUE),
            (0.75, 3.15, "CPU 内存（bounce）", CORAL),
            (0.75, 4.30, "网卡 NIC", TEAL),
        ]
        for x, y, name, col in nodes:
            box(s, x, y, 5.5, 0.72, BG_CARD2, col, 0.08)
            tb(s, x, y, 5.5, 0.72, name, 18, WHITE, True, align="center",
               anchor=MSO_ANCHOR.MIDDLE)
        tb(s, 0.75, 5.15, 5.5, 1.3,
           "两次 PCIe 拷贝，吃 CPU 带宽。\n程序看起来「先 cudaMemcpy 再 send」。",
           14, MUTED)
        # right: gdr
        box(s, 6.75, 1.22, 6.05, 5.55, BG_CARD, LINE)
        tb(s, 6.97, 1.38, 5.6, 0.36, "GDR 路径：网卡直读 GPU", 18, GREEN, True)
        box(s, 7.05, 2.15, 5.5, 0.72, BG_CARD2, BLUE, 0.08)
        tb(s, 7.05, 2.15, 5.5, 0.72, "GPU 显存", 18, WHITE, True, align="center",
           anchor=MSO_ANCHOR.MIDDLE)
        box(s, 7.05, 3.85, 5.5, 0.72, BG_CARD2, TEAL, 0.08)
        tb(s, 7.05, 3.85, 5.5, 0.72, "网卡 NIC 直接 DMA", 18, WHITE, True, align="center",
           anchor=MSO_ANCHOR.MIDDLE)
        tb(s, 8.9, 3.00, 2.0, 0.7, "↓  少一跳", 16, GREEN, True, align="center")
        tb(s, 7.05, 4.85, 5.5, 1.55,
           "这就是 GPUDirect RDMA。\nGDRCopy 的 gdr_pin_buffer()，就是在问内核：\n「这块 GPU 地址，能不能给网卡当 DMA 源？」",
           14, SOFT)
        add_notes(s, """
用 90 秒讲图，不要展开 PCIe TLP。
点一句：我们后面测的 write_lat_gpu --gpu-mode=peermem，走的就是右边这条。
GDRCopy 不是「拷贝库」那么简单，它的第一步 pin，是在内核里找「这 VA 对应哪段 GPU 物理页」。找不到就直接失败，后面的 BAR mapping 根本不会发生。
""")

    # ---------- 5 hotel ----------
    def s_hotel():
        s = add()
        header(s, "01  背景", "贯穿全场的比喻：酒店前台的访客登记本", len(slides), TOTAL[0])
        rows = [
            ("GPU 显存", "酒店房间", "物理上一直在，4090 的 BAR1 也能被别人看到。", BLUE),
            ("CUDA VA", "房卡上的门牌", "应用程序手里的地址，例如 0x304200000。", TEAL),
            ("hMemory", "前台内部档案号", "NVIDIA RM 给这段显存的对象句柄，每次运行都不同。", PURPLE),
            ("0x503c 登记本", "「允许网卡来访」名单", "ThirdPartyP2P。Tesla 的 libcuda 会自己建；GeForce 不会。", AMBER),
            ("nvidia_p2p_get_pages", "网卡保安来查名单", "名单上没有这个门牌 → 0x57 查无此人 → pin 返回 -22。", CORAL),
        ]
        y = 1.20
        for left, mid, right, col in rows:
            box(s, 0.48, y, 12.35, 1.08, BG_CARD, LINE, 0.05)
            rect(s, 0.48, y, 0.10, 1.08, col)
            tb(s, 0.75, y + 0.12, 2.7, 0.84, left, 15, WHITE, True, mono=True,
               anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 3.5, y + 0.12, 2.6, 0.84, mid, 16, col, True, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 6.2, y + 0.16, 6.4, 0.80, right, 15, SOFT, anchor=MSO_ANCHOR.MIDDLE)
            y += 1.16
        add_notes(s, """
这是全场最重要的一页。后面说「建登记本」「写门牌」「查无此人」，都指这一张表。
强调 handle 每次不同：所以我们的 .so 里没有写死 GDB 里看到的地址，写死的是 RM 公开协议（class / cmd），房间号运行时从 ioctl 里抓。
0x503c 就是 NVIDIA 给第三方 P2P 用的对象类名 NV50_THIRD_PARTY_P2P。听众不必背，记住「登记本」即可。
""")

    # ---------- 6 our machine ----------
    def s_machine():
        s = add()
        header(s, "01  背景", "我们机器上，已经通了什么、还缺哪扇门", len(slides), TOTAL[0])
        card_title(s, 0.48, 1.22, 6.05, 5.55, "已经有的（不要再动）",
                   "•  aikitoria 575.64.05-p2p：GPU↔GPU BAR1 P2P\n"
                   "•  BAR1 扩到 32 GB\n"
                   "•  驱动 575.64.05，CUDA 12.9，AArch64\n"
                   "•  卡是 RTX 4090 D\n\n"
                   "这只证明：别人可以从 BAR1 碰到这块 GPU 内存。\n"
                   "它不自动等于「网卡可以按 CUDA VA 来 pin」。",
                   GREEN, 18, 16)
        card_title(s, 6.75, 1.22, 6.05, 5.55, "还没有的（今天的题目）",
                   "•  GDRCopy：gdr_pin_buffer()\n"
                   "•  内核符号：nvidia_p2p_get_pages()\n"
                   "•  URMA：write_lat_gpu --gpu-mode=peermem\n\n"
                   "原则：不要为了 GDR 再去改 4090 的 NVIDIA 内核。\n"
                   "Harry 用 5090 证明了——最后一道闸在用户态。",
                   AMBER, 18, 16)
        add_notes(s, """
分清两件已被同事搞混的事：
1) GPU 跟 GPU 互访（P2P）—— aikitoria 已经做了，保持原样。
2) 网卡按 CUDA 虚拟地址去 pin GPU 页 —— 这是 GDRCopy / peermem，今天的主题。
当场再说一次：禁止再给 4090 打 NVIDIA 内核 GDR 补丁。内核那条路 Harry 试过，不是最后闸门；我们还需要的是 legacy 登记本，内核里改 dma-buf 能力位也补不上。
""")

    # ---------- 7 error ----------
    def s_error():
        s = add()
        header(s, "01  背景", "失败长什么样：不是网卡坏了，是名单上没你", len(slides), TOTAL[0])
        box(s, 0.48, 1.22, 12.35, 2.15, BG_CODE, RGBColor(0x5A, 0x2A, 0x32))
        tb(s, 0.75, 1.40, 11.8, 1.8,
           [
               ("gdr_pin_buffer(...)  =  -22     /* EINVAL */", 20, CORAL, True, True),
               ("NVIDIA RM status     =  0x57    /* NV_ERR_OBJECT_NOT_FOUND */", 20, CORAL, True, True),
               ("人话：保安来查登记本，这个 CUDA VA 根本不在本上。", 16, SOFT, False, False),
           ])
        items = [
            ("会的（Tesla / 数据中心卡）", GREEN,
             "libcuda 自己报 attr 116=1，并且自己 Alloc 0x503c，再 REGISTER_VA_SPACE / REGISTER_VIDMEM。"),
            ("不会的（GeForce / 4090）", CORAL,
             "libcuda 报 attr 116=0，也从不发那三步。内核里那本登记册是空的，甚至根本不存在。"),
            ("所以缺口是软件政策", AMBER,
             "不是「4090 的 BAR 不能被 DMA」。BAR1 P2P 已经证明物理通道在。缺的是「允许第三方按 VA 查找」这份手续。"),
        ]
        for i, (t, c, b) in enumerate(items):
            card_title(s, 0.48 + i * 4.2, 3.60, 4.05, 3.15, t, b, c, 16, 14)
        add_notes(s, """
把 -22 和 0x57 写在黑板上。后面日志里只要再出现这两个数，听众立刻知道是登记本问题。
Tesla 对照非常重要：同一份 GDRCopy，在 V100 上不用我们的 .so 就能 pin。说明 GDRCopy 没写错，是 GeForce 的 libcuda 少做了功课。
0x57 的官方名字是 NV_ERR_OBJECT_NOT_FOUND。对象不是 GPU，是「这个 VA 对应的第三方 P2P 登记项」。
""")

    # ---------- 8 harry 1 ----------
    def s_harry_start():
        s = add()
        header(s, "02  Harry 的 5090", "他的起点：内核 P2P 有了，CUDA 仍然说「不支持」", len(slides), TOTAL[0])
        tb(s, 0.48, 1.18, 12.3, 0.4,
           "文章：harrychen.xyz/2026/05/20/enable-gpudirect-rdma-on-rtx-5090/",
           14, BLUE, True, mono=True)
        steps = [
            ("1", "已有 patched 内核 P2P", "5090 上 GPU 互访已经能做，和我们 aikitoria BAR1 很像。"),
            ("2", "CUDA 能力位仍是 0", "GDR / dma-buf 相关 attribute 返回 0。"),
            ("3", "导出句柄直接 801", "cuMemGetHandleForAddressRange = CUDA_ERROR_NOT_SUPPORTED。"),
            ("4", "先去改内核能力位", "强迫 dma_buf_supported、DMABUF_CAPABILITY_YES。"),
            ("5", "CUDA 还是报 0", "结论：内核点头，用户态库可以继续摇头。"),
        ]
        y = 1.65
        for n, t, b in steps:
            box(s, 0.48, y, 12.35, 0.95, BG_CARD, LINE, 0.05)
            box(s, 0.65, y + 0.22, 0.52, 0.52, AMBER, radius=0.15)
            tb(s, 0.65, y + 0.22, 0.52, 0.52, n, 18, BG, True, align="center",
               anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 1.35, y + 0.12, 4.4, 0.72, t, 17, WHITE, True, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 5.8, y + 0.16, 6.8, 0.68, b, 15, SOFT, anchor=MSO_ANCHOR.MIDDLE)
            y += 1.02
        add_notes(s, """
用 Harry 的时间线当「我们为什么不再打内核」的外部证据。
他比我们更早把「改内核能力位」做了一遍：NV2080_CTRL_GPU_INFO_INDEX_DMABUF_CAPABILITY_YES 之类。CUDA 用户态照样报 0。
给听众的翻译：驱动里的「我支持」和 libcuda 里的「我宣布支持」是两套旗。应用只信 libcuda。
我们 4090 上如果继续打 NVIDIA 内核 GDR 补丁，会重复他第 4、第 5 步。
""")

    # ---------- 9 harry bit ----------
    def s_harry_bit():
        s = add()
        header(s, "02  Harry 的 5090", "闸门在 libcuda：属性 116 看的是一个软件 bit", len(slides), TOTAL[0])
        card_title(s, 0.48, 1.22, 6.05, 3.35, "他反汇编 cuDeviceGetAttribute",
                   "属性 116（GPU_DIRECT_RDMA_SUPPORTED）在 libcuda 里并不是去问内核「你到底行不行」。\n\n"
                   "它看两样东西：\n"
                   "•  设备对象上某个字节的 0x20 位\n"
                   "•  一个全局的 gdr_enabled 开关\n\n"
                   "两位都过，用户态才对外说「我支持 GDR」。",
                   BLUE, 17, 15)
        card_title(s, 6.75, 1.22, 6.05, 3.35, "进程里把 bit 拨开，立刻变 1",
                   "他没有重编驱动。只在正在跑的进程里把那个 0x20 拨开。\n\n"
                   "cuDeviceGetAttribute(116) 马上变成 1，dma-buf 导出也开始工作。\n\n"
                   "这是全篇文章最关键的实验：证明闸门在用户态策略，不在 5090 硬件。",
                   GREEN, 17, 15)
        box(s, 0.48, 4.75, 12.35, 2.05, BG_CODE, LINE)
        tb(s, 0.7, 4.90, 12.0, 1.75,
           [
               ("libcuda.so 设备初始化里原来是：     or  $0x40, %edx", 16, SOFT, False, True),
               ("他改成：                           or  $0x60, %edx    # 0x40 | 0x20", 16, GREEN, True, True),
               ("文件偏移（590.48.01）:  0x42d69b   一个字节  0x40 -> 0x60", 16, AMBER, False, True),
           ])
        add_notes(s, """
慢慢讲。基础程序员能懂「函数返回值被一个 if 挡住了」。
Harry 的手法是：找到挡住 116 的那个 bit，在 libcuda 拷贝里改一字节，用 LD_LIBRARY_PATH 加载这份拷贝。
他中间弃过一条不完整的 shim（cuGetProcAddress 没转发全）。最终方案是改文件，不是 LD_PRELOAD。
我们后面会说：这个 bit 我们也要「变成 1」，但我们不改 libcuda 文件，而是挟持 cuDeviceGetAttribute，直接把返回值改写成 1。效果类似，手段更干净。
""")

    # ---------- 10 harry result ----------
    def s_harry_result():
        s = add()
        header(s, "02  Harry 的 5090", "他打开的是 dma-buf 大门，作者写明：不是旧 peermem", len(slides), TOTAL[0])
        box(s, 0.48, 1.22, 12.35, 1.55, BG_CARD, GREEN)
        tb(s, 0.75, 1.40, 11.8, 1.25,
           "NCCL allreduce busbw 大约 8.87 → 19.93 GB/s，日志出现 NET/IB/.../GDRDMA。\n"
           "对 5090 + NCCL 这条线，一字节足够，而且效果是真的。",
           18, WHITE)
        card_title(s, 0.48, 3.00, 6.05, 3.75, "Harry 打开的",
                   "•  CUDA 对外宣布支持 GDR\n"
                   "•  dma-buf 导出\n"
                   "•  新一代 NCCL / IB 走 GDRDMA\n\n"
                   "应用不再需要 nvidia-peermem.ko 那条老路径。",
                   GREEN, 18, 16)
        card_title(s, 6.75, 3.00, 6.05, 3.75, "Harry 自己说没打开的",
                   "•  旧的 nvidia-peermem\n"
                   "•  nvidia_p2p_get_pages 那套查找\n"
                   "•  经典 GDRCopy pin\n\n"
                   "文章里写得很清楚：不要以为改完这一字节，老 peermem 就活了。",
                   CORAL, 18, 16)
        add_notes(s, """
这是「借鉴」和「照抄」的分界线。必须停一停。
Harry 成功了，成功的是 NCCL + dma-buf。
我们 4090 上要跑的是 GDRCopy 和 URMA write_lat_gpu --gpu-mode=peermem。这条线内核走 nvidia_p2p_get_pages，查的是 ThirdPartyP2P 登记本。
所以：Harry 证明了「别再打内核」；Harry 没有替我们写登记本。下一页说我们为什么不能只抄一字节。
""")

    # ---------- 11 why not one byte ----------
    def s_not_one_byte():
        s = add()
        header(s, "03  借鉴", "为什么我们不能只抄「改一字节」", len(slides), TOTAL[0])
        bullets(s, 0.48, 1.20, 12.3, 2.2, [
            "GDRCopy 的 check_gdr_support() 会读属性 116。只改 bit / 只撒谎返回 1，它肯开始 pin。",
            "pin 的下一步是 nvidia_p2p_get_pages(va)。这一步不看属性 116，看 RM 里有没有登记项。",
            "GeForce 的 libcuda 从来不发 Alloc 0x503c，也从来不发 REGISTER_VIDMEM。登记本是空的。",
            "空本 + 正确的 116 = 仍然 0x57。我们在 4090 上就踩过这个坑。",
        ], 17)
        box(s, 0.48, 3.55, 12.35, 3.20, BG_CARD, AMBER)
        tb(s, 0.75, 3.75, 11.8, 0.4, "一句话", 16, AMBER, True)
        tb(s, 0.75, 4.25, 11.8, 2.15,
           "改 bit 是打开「我宣布支持」这扇门；\n"
           "写登记本是打开「保安查得到你」那扇门。\n"
           "Harry 的 5090 只要第一扇（dma-buf）。我们的 4090 两扇都要，而且第二扇才是 GDRCopy 的生死线。",
           20, WHITE)
        add_notes(s, """
可以问听众：如果只 hook cuDeviceGetAttribute 返回 1，会发生什么？
答案：GDRCopy 通过能力检查，然后 pin 失败，错误还是 -22。日志里不会有 REGISTER_VIDMEM。
这就是我们早期尝试的真实形态。所以必须继续往下做 ioctl 和 VMM。
""")

    # ---------- 12 our failures ----------
    def s_failures():
        s = add()
        header(s, "03  借鉴", "我们走过的弯路（按时间）", len(slides), TOTAL[0])
        rows = [
            ("1", "继续打 NVIDIA 内核", "以为缺的是 dma-buf / P2P 能力位。", "Harry 已经证明内核不是最后闸门。我们决定停。"),
            ("2", "只骗属性 116 / 110", "GDRCopy 肯进 pin。", "get_pages 仍 0x57。只打开了「宣布支持」。"),
            ("3", "普通 cuMemAlloc 后去登记", "VA 有了，以为能写本。", "RM 里常常没有可登记的 class 0x40 hMemory。"),
            ("4", "只 hook ioctl，不 hook dlsym", "命令行小工具好像好了。", "urma_perftest 用 dlopen+dlsym，绕过 PLT，钩子没套上。"),
            ("5", "全局 LD_PRELOAD / sudo", "图省事。", "sudo 丢掉环境变量；不相关进程也会加载 .so。"),
        ]
        y = 1.18
        for n, a, b, c in rows:
            box(s, 0.48, y, 12.35, 1.08, BG_CARD, LINE, 0.04)
            box(s, 0.62, y + 0.28, 0.50, 0.50, CORAL, radius=0.14)
            tb(s, 0.62, y + 0.28, 0.50, 0.50, n, 16, WHITE, True, align="center",
               anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 1.28, y + 0.14, 3.3, 0.80, a, 15, WHITE, True, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 4.65, y + 0.14, 3.7, 0.80, b, 14, MUTED, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 8.45, y + 0.14, 4.2, 0.80, c, 14, SOFT, anchor=MSO_ANCHOR.MIDDLE)
            y += 1.14
        add_notes(s, """
用「我们自己也栽过」建立信任，不是炫技。
第 3 条要多停 10 秒：这是后面为什么必须把 cuMemAlloc 换成 VMM 的伏笔。普通分配走 UVM，RM 里没有一张能拿去 REGISTER_VIDMEM 的 class 0x40 档案。
第 4 条对应代码末尾的 dlsym()。第 5 条对应 isolate/run.sh。后面都会回到这两条。
""")

    # ---------- 13 borrow ----------
    def s_borrow():
        s = add()
        header(s, "03  借鉴", "借鉴什么、不借鉴什么", len(slides), TOTAL[0])
        card_title(s, 0.48, 1.22, 4.05, 5.55, "从 Harry 借",
                   "•  最后闸门在 userspace\n"
                   "•  不要再为 GDR 改内核\n"
                   "•  属性 116 是软件旗，可以拨\n\n"
                   "不借：一字节 patch libcuda 文件。\n"
                   "不借：只做 dma-buf / NCCL 这条线当全部答案。",
                   AMBER, 17, 15)
        card_title(s, 4.68, 1.22, 4.05, 5.55, "从公开 writeup 借",
                   "remote-lab.net\n「Enabling GPUDirect RDMA on GeForce」\n\n"
                   "github.com/mcornea/\ngeforce-gpudirect-rdma\n\n"
                   "借的是：ioctl 抓 RM handle，自己建 NV503C，再 REGISTER_VIDMEM。这是登记本路线。",
                   BLUE, 16, 14)
        card_title(s, 8.88, 1.22, 3.95, 5.55, "我们补上的",
                   "•  cuMemAlloc → VMM\n"
                   "•  挟持 ioctl 只拆 RM_ALLOC\n"
                   "•  挟持 dlsym（UMDK）\n"
                   "•  吞掉 SYNC_MEMOPS 801\n"
                   "•  isolate + run.sh\n"
                   "•  不改 libcuda 文件\n"
                   "•  不全局 preload",
                   GREEN, 17, 15)
        add_notes(s, """
三列一起指：思想来自 Harry，手续来自 remote-lab/mcornea，工程落地是我们的 gdr_geforce_hook.c。
对听众说：开源世界上这条路不是我们发明的协议，NV503C 是 NVIDIA 公开类。我们做的是把手续自动化，并解决 4090 + UMDK 上的几个坑（VMM、dlsym、801、隔离运行）。
""")

    # ---------- 14 birdview ----------
    def s_bird():
        s = add()
        header(s, "04  原理", "整条链路鸟瞰：.so 插在哪里", len(slides), TOTAL[0])
        labels = [
            ("应用\nGDRCopy / URMA", CORAL),
            ("我们的 .so\n插队的前台", AMBER),
            ("真 libcuda\n仍做大部分事", BLUE),
            ("ioctl\n递工单", PURPLE),
            ("NVIDIA RM\n内核前台", TEAL),
        ]
        for i, (name, col) in enumerate(labels):
            x = 0.40 + i * 2.55
            box(s, x, 1.40, 2.35, 1.55, BG_CARD, col, 0.08)
            tb(s, x + 0.08, 1.48, 2.2, 1.4, name, 15, WHITE, True, align="center",
               anchor=MSO_ANCHOR.MIDDLE)
            if i < 4:
                arrow(s, x + 2.30, 1.95, 0.32, 0.28, col)
        box(s, 0.48, 3.20, 12.35, 3.55, BG_CARD, LINE)
        tb(s, 0.75, 3.38, 11.8, 0.4, "我们的 .so 在这一跳里做四件事", 16, AMBER, True)
        jobs = [
            ("骗", "cuDeviceGetAttribute(116/110) 强制返回 1"),
            ("换", "cuMemAlloc 改成 VMM，逼 RM 留下 class 0x40"),
            ("看", "ioctl 成功后再拆 RM_ALLOC，记下房间号"),
            ("补", "自己 Alloc 0x503c，REGISTER_VA_SPACE + VIDMEM"),
        ]
        for i, (k, v) in enumerate(jobs):
            x = 0.75 + (i % 2) * 6.0
            y = 3.90 + (i // 2) * 1.25
            tb(s, x, y, 0.7, 0.9, k, 22, AMBER, True, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, x + 0.8, y, 5.0, 0.9, v, 16, SOFT, anchor=MSO_ANCHOR.MIDDLE)
        add_notes(s, """
用「骗 / 换 / 看 / 补」四个汉字当记忆钩。后面四段就是这四个字的展开。
强调：真 libcuda 还在，内核还在。我们不是重写 CUDA，是补它不肯做的手续。
箭头从左到右是一次分配+pin 的调用方向。返回时 get_pages 从最右边的 RM 里已经能查到登记项。
""")

    # ---------- 15 ldpreload ----------
    def s_preload():
        s = add()
        header(s, "04  原理", "LD_PRELOAD：在门口插队，不是改文件", len(slides), TOTAL[0])
        code_card(s, 0.48, 1.20, 6.2, 2.35, [
            "export LD_PRELOAD=$PWD/libgdr_geforce_hook.so",
            "./gdrcopy_copylat",
            "",
            "# 动态链接器 ld.so 的加载顺序：",
            "#   1. 我们的 .so",
            "#   2. 可执行文件依赖的 libc / libcuda",
        ], "shell")
        box(s, 6.90, 1.20, 5.9, 2.35, BG_CARD, LINE)
        tb(s, 7.1, 1.35, 5.5, 2.05,
           "解析未定义符号时，从先加载的库开始搜。\n"
           "我们导出了和 libc / libcuda 同名的函数：\n"
           "ioctl、cuMemAlloc_v2、cuDeviceGetAttribute…\n"
           "于是 call ioctl / bl cuMemAlloc 先命中我们。",
           15, SOFT)
        rows = [
            ("是", GREEN, "符号拦截。进程里的函数指针指向我们。"),
            ("不是", CORAL, "改 libcuda.so 某个文件偏移（Harry 的一字节）。"),
            ("不是", CORAL, "patch 内核、加系统调用。"),
            ("自己还要调真货", AMBER, "dlsym(RTLD_NEXT, \"ioctl\")，从我们后面继续搜。千万不要再调自己，否则递归死循环。"),
        ]
        y = 3.75
        for a, c, b in rows:
            box(s, 0.48, y, 12.35, 0.80, BG_CARD, LINE, 0.04)
            tb(s, 0.7, y, 2.6, 0.80, a, 16, c, True, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 3.4, y, 9.2, 0.80, b, 16, SOFT, anchor=MSO_ANCHOR.MIDDLE)
            y += 0.85
        add_notes(s, """
给没写过动态链接的人：把 LD_PRELOAD 想成「前台先排一个自己人」。客人喊「ioctl」，自己人先接，再决定转给真正的前台。
对比 Harry：他改的是 libcuda 文件里的一条 or 指令；我们连文件都不动。升级驱动时不用重算偏移。
RTLD_NEXT 一定要讲：否则听众会问「你 hook 了 ioctl，自己怎么发 ioctl？」答案在 call_ioctl()。
绑不住的情况也提一句：对方如果 syscall(SYS_ioctl) 或静态链接 CUDA，我们就套不住。575 + GDRCopy 走的是动态符号，所以有效。
""")

    # ---------- 16 dlsym ----------
    def s_dlsym():
        s = add()
        header(s, "04  原理", "只插 PLT 不够：urma_perftest 用 dlsym 绕过门口", len(slides), TOTAL[0])
        card_title(s, 0.48, 1.22, 6.05, 3.45, "普通程序",
                   "编译时记下 cuMemAlloc_v2 这个名字。\n运行时走 PLT / GOT，动态链接器帮你找到符号。\nLD_PRELOAD 对这条路天生有效。",
                   BLUE, 18, 16)
        card_title(s, 6.75, 1.22, 6.05, 3.45, "UMDK 的 loader",
                   "自己 dlopen(\"libcuda.so.1\")，\n再 dlsym(handle, \"cuMemAlloc_v2\")。\n拿到的是 libcuda 里的真函数指针，不经过 PLT。\n只 hook 函数本身，套不住它。",
                   CORAL, 18, 16)
        code_card(s, 0.48, 4.85, 12.35, 2.00, [
            "void *dlsym(void *handle, const char *name) {",
            "    if (handle != RTLD_NEXT && is_our_name(name))",
            "        return (void *)our_function;   // 把钩子还回去",
            "    return real_dlsym(handle, name);   // RTLD_NEXT 原样，才能拿到真货",
            "}",
        ], "gdr_geforce_hook.c  ·  dlsym")
        add_notes(s, """
这是「基础程序员」最容易漏的一页。很多人以为 LD_PRELOAD 万能。
现场演示口播：urma_perftest 的 cuda-loader 就是 dlopen + dlsym。所以我们连 libc 的 dlsym 也拦截。
内部查真符号一律走 real_dlsym / dlvsym，禁止走我们自己导出的 dlsym，否则一查自己就套娃。
handle == RTLD_NEXT 时不要替换，否则 lookup_next() 会拿到我们自己。
""")

    # ---------- 17 ioctl what ----------
    def s_ioctl_what():
        s = add()
        header(s, "05  挟持 ioctl", "ioctl 是什么：给驱动递一张写好的工单", len(slides), TOTAL[0])
        tb(s, 0.48, 1.18, 12.3, 0.7,
           "你已经会的：open(\"/dev/xxx\") 得到 fd。ioctl 就是拿着这个 fd，再附一张结构体，请内核做一件设备特有的事。",
           16, SOFT)
        # envelope
        box(s, 0.48, 2.00, 12.35, 2.55, BG_CARD, LINE)
        tb(s, 0.7, 2.12, 12.0, 0.35, "一张 NVIDIA 工单长这样", 16, AMBER, True)
        fields = [
            ("magic = 'F'", "这是 NVIDIA 的信封口令"),
            ("nr = 0x2B", "RM_ALLOC：请新建一个对象"),
            ("nr = 0x2A", "RM_CONTROL：请对已有对象发命令"),
            ("结构体里的 class", "要建哪一类房间（设备 / 显存 / 登记本）"),
            ("返回的 handle", "房间号，写回同一张工单"),
        ]
        for i, (a, b) in enumerate(fields):
            x = 0.7 + i * 2.4
            tb(s, x, 2.55, 2.25, 0.7, a, 13, WHITE, True, mono=True)
            tb(s, x, 3.20, 2.25, 1.1, b, 13, MUTED)
        card_title(s, 0.48, 4.75, 6.05, 2.05, "听众只需要记住",
                   "CUDA 并没有魔法。它跟内核说话，靠的就是反复 ioctl。\n我们挟持的是 libc 的 ioctl()，所以能看见这些工单。",
                   BLUE, 16, 15)
        card_title(s, 6.75, 4.75, 6.05, 2.05, "我们不挟持的",
                   "socket、tty、eventfd、framebuffer 的 ioctl。\n形状不像 RM_ALLOC 的，看一眼类型就原样转发。",
                   GREEN, 16, 15)
        add_notes(s, """
用「工单 / 信封」把 ioctl 从内核黑话里救出来。
magic 'F' + nr 是 _IOWR 宏编出来的。听众不必会算，知道「我们先看这是不是 NVIDIA 建对象的信封」即可。
0x2B / 0x2A / 0x28 对应代码里的 NV_ESC_RM_ALLOC、RM_CONTROL、RM_ALLOC_OBJECT。
强调「返回的 handle 写在同一张工单里」：所以我们是在真 ioctl 成功之后再读这张纸，不是拦截并改写请求。
""")

    # ---------- 18 ioctl discipline ----------
    def s_ioctl_code():
        s = add()
        header(s, "05  挟持 ioctl", "纪律：先让真 ioctl 做完，成功了再偷看", len(slides), TOTAL[0])
        code_card(s, 0.48, 1.18, 12.35, 3.55, [
            "int ioctl(int fd, unsigned long request, ...) {",
            "    void *arg = va_arg(...);",
            "    if (!is_rm_alloc_ioctl(request))",
            "        return call_ioctl(fd, request, arg);   // 不是建对象工单：立刻转发",
            "",
            "    int ret = call_ioctl(fd, request, arg);    // 先做真事",
            "    if (ret != 0 || !is_nvidia_rm_fd(fd))",
            "        return ret;                            // 失败或不是 NVIDIA 窗口：不拆包",
            "",
            "    after_rm_alloc(fd, request, arg);          // 成功后再记账",
            "    return ret;",
            "}",
        ], "gdr_geforce_hook.c  ·  ioctl()   约 540 行")
        pills = [
            ("不改工单内容", "CUDA 以为自己在跟内核说话"),
            ("不挡路", "我们是事后抄房间号的人"),
            ("自己发工单走 call_ioctl", "不会再次进这个函数"),
        ]
        for i, (t, b) in enumerate(pills):
            box(s, 0.48 + i * 4.2, 4.95, 4.05, 1.75, BG_CARD, LINE)
            tb(s, 0.68 + i * 4.2, 5.10, 3.7, 0.45, t, 16, AMBER, True)
            tb(s, 0.68 + i * 4.2, 5.55, 3.7, 0.95, b, 15, SOFT)
        add_notes(s, """
对着代码念，不要跳。这是讲座的核心页。
三步：形状过滤 → 真调用 → 认人再拆包。
call_ioctl 内部是 real_ioctl = dlsym(RTLD_NEXT, \"ioctl\")。我们自己补登记本时也走它，注释里写了「不进这里」。
is_nvidia_rm_fd 放在成功之后，是为了：非 NVIDIA 的、碰巧也用 'F' 的 ioctl（framebuffer FBIO*）不要被我们 fstat、加锁。先做完、再确认窗口。
""")

    # ---------- 19 who ----------
    def s_who():
        s = add()
        header(s, "05  挟持 ioctl", "认人：只拆 NVIDIA 窗口上的信封", len(slides), TOTAL[0])
        card_title(s, 0.48, 1.22, 6.05, 5.55, "怎样认出 NVIDIA RM",
                   "1. 字符设备，major = 195\n"
                   "2. 或者路径是 /dev/nvidiactl、/dev/nvidia0..N\n"
                   "3. /proc/self/fd/N 的软链对得上\n\n"
                   "认不出就不拆包。socket、终端、eventfd、空设备，全部当透明管道。",
                   TEAL, 17, 16)
        code_card(s, 6.75, 1.22, 6.05, 5.55, [
            "static int is_rm_alloc_ioctl(request) {",
            "  type == 'F'",
            "  nr == 0x2B && size 对得上  // ALLOC",
            "   || nr == 0x28            // ALLOC_OBJECT",
            "}",
            "",
            "is_nvidia_rm_fd(fd):",
            "  fstat -> S_ISCHR",
            "  major(st.st_rdev) == 195",
            "  或 readlink(/proc/self/fd/N)",
            "     是 /dev/nvidiactl|nvidiaN",
        ], "过滤条件")
        add_notes(s, """
有人会问：LD_PRELOAD 之后每个进程的 ioctl 不都被换了吗？
答：符号层面是的，函数体不是。非 RM_ALLOC 形状的请求，判断 type/nr 之后立刻 call_ioctl，不加锁、不 fstat。
这就是为什么文档写「即使误进登录环境，cat / vim 也不会被我们拆包」。尽管如此，我们仍然要求只用 isolate/run.sh，不要全局 preload。
""")

    # ---------- 20 family ----------
    def s_family():
        s = add()
        header(s, "05  挟持 ioctl", "RM 对象是家谱：我们要抄的是这几代房间号", len(slides), TOTAL[0])
        nodes = [
            (5.15, 1.25, 3.0, "hClient  根", "一次 CUDA 上下文的客户", AMBER),
            (2.05, 2.85, 3.0, "hDevice  0x80", "一块 GPU", BLUE),
            (8.25, 2.85, 3.0, "hSubdevice  0x2080", "0x503c 必须挂在它下面", TEAL),
            (2.05, 4.55, 3.0, "hVASpace  0x90f1", "这块 GPU 的虚拟地址空间", PURPLE),
            (8.25, 4.55, 3.0, "hMemory  0x40", "这一段可登记的显存档案", GREEN),
        ]
        for x, y, w, title, sub, col in nodes:
            box(s, x, y, w, 1.15, BG_CARD, col, 0.08)
            tb(s, x + 0.1, y + 0.12, w - 0.2, 0.45, title, 16, WHITE, True, mono=True,
               align="center")
            tb(s, x + 0.1, y + 0.55, w - 0.2, 0.48, sub, 13, MUTED, align="center")
        # simple connectors as thin rects
        rect(s, 6.55, 2.40, 0.06, 0.45, LINE)
        rect(s, 3.50, 2.40, 6.30, 0.06, LINE)
        rect(s, 3.50, 4.00, 0.06, 0.55, LINE)
        rect(s, 9.70, 4.00, 0.06, 0.55, LINE)
        tb(s, 0.48, 5.90, 12.3, 1.15,
           "class 是房间类型，写死在协议里。handle 是房间号，每次入住都不同。\n"
           "所以 .so 里没有写死你 GDB 里看到的 0x5c00009d，只写死 0x80 / 0x2080 / 0x90f1 / 0x40 / 0x503c。",
           16, SOFT)
        add_notes(s, """
画家谱比念结构体有效。指着说：登记本 0x503c 的父亲必须是 subdevice；REGISTER_VIDMEM 要的是 class 0x40 的 hMemory，加上这块 vaspace 里的 VA。
0x40 的正式名是 NV01_MEMORY_LOCAL_USER。听众记住「能拿去登记的显存档案」即可。
""")

    # ---------- 21 track ----------
    def s_track():
        s = add()
        header(s, "05  挟持 ioctl", "track_alloc：只记账，不改 CUDA 正在做的事", len(slides), TOTAL[0])
        code_card(s, 0.48, 1.18, 12.35, 3.70, [
            "if (hClass == 0x80)   记下 hClient / hDevice     // DEVICE",
            "if (hClass == 0x2080) 记下 hSubdevice            // 挂在目标 GPU 下",
            "if (hClass == 0x90f1) 记下 hVASpace, handles_ready = 1",
            "if (capturing_memory && 是显存 class)",
            "    captured_hMemory = hNew;   // 开窗口期间出现的那张档案",
        ], "track_alloc()  ·  约 333 行")
        bullets(s, 0.55, 5.10, 12.2, 1.85, [
            "status 必须是 0，handle 必须非 0，否则当这次没建成。",
            "换了一块 GPU / 换了一个 client，就把旧房间号清掉，避免张冠李戴。",
            "环境变量 GPUDIRECT_GPU 选第几块卡，默认 0。",
        ], 16)
        add_notes(s, """
对着家谱讲这一页。ioctl 成功返回后，工单里的 hObjectNew 就是新房间号，hClass 是类型。
capturing_memory 是后面 VMM 那一节的开关：只在 cuMemCreate 前后的窗口里抓 hMemory，免得把 CUDA 自己其它零碎分配当成「这一块要登记的显存」。
""")

    # ---------- 22 why vmm ----------
    def s_why_vmm():
        s = add()
        header(s, "06  换分配", "为什么普通 cuMemAlloc 写不进登记本", len(slides), TOTAL[0])
        card_title(s, 0.48, 1.22, 6.05, 5.55, "cuMemAlloc 常见路径",
                   "走 UVM。用户态拿到一个 VA。\n"
                   "RM 里常常没有一张能拿去 REGISTER_VIDMEM 的 class 0x40 档案。\n\n"
                   "就像：你有房卡门牌，前台档案柜里却没有对应的纸质档案。保安查的是档案柜。",
                   CORAL, 18, 16)
        card_title(s, 6.75, 1.22, 6.05, 5.55, "VMM 路径（我们换成这个）",
                   "cuMemCreate 会让 RM 真的 Alloc 一个 class 0x40 对象。\n"
                   "ioctl 钩子在这个窗口里能抓到 hMemory。\n"
                   "再 cuMemMap 到一段自己 reserve 的 VA。\n\n"
                   "门牌和档案对得上，才能写进登记本。",
                   GREEN, 18, 16)
        add_notes(s, """
这是弯路第 3 条的展开。现场可以说：我们早期对着普通 cuMemAlloc 出来的 VA 发 REGISTER_VIDMEM，日志是 no hMemory。
UVM 是 CUDA 的统一虚拟内存。听众只需知道「这条路径不保证 RM 里有可登记对象」。
VMM = Virtual Memory Management API：cuMemCreate / cuMemMap / cuMemSetAccess。数据中心卡上 GDR 也更推荐这条。
""")

    def s_vmm_steps():
        s = add()
        header(s, "06  换分配", "cuMemAlloc 被换成这四步", len(slides), TOTAL[0])
        steps = [
            ("1", "Reserve", "cuMemAddressReserve 先占一段 64K 对齐的 GPU VA。"),
            ("2", "Create", "cuMemCreate 向 RM 要物理 FB。这里打开抓 hMemory 窗口。"),
            ("3", "Map + Access", "cuMemMap 把门牌贴到档案上，再 SetAccess 允许读写。"),
            ("4", "Register", "拿 (hMemory, VA, size) 写进 0x503c 登记本。"),
        ]
        for i, (n, t, b) in enumerate(steps):
            x = 0.48 + i * 3.2
            box(s, x, 1.25, 3.05, 3.55, BG_CARD, LINE, 0.07)
            box(s, x + 0.20, 1.45, 0.55, 0.55, AMBER, radius=0.14)
            tb(s, x + 0.20, 1.45, 0.55, 0.55, n, 18, BG, True, align="center",
               anchor=MSO_ANCHOR.MIDDLE)
            tb(s, x + 0.18, 2.20, 2.7, 0.55, t, 20, WHITE, True)
            tb(s, x + 0.18, 2.85, 2.7, 1.7, b, 15, SOFT)
            if i < 3:
                arrow(s, x + 2.95, 2.70, 0.32, 0.28, AMBER)
        tb(s, 0.48, 5.00, 12.3, 1.85,
           "GeForce 上 cuMemCreate 带 gpuDirectRDMACapable=1 常常返回 101。\n"
           "我们丢掉这面旗再试一次——物理页仍然能建出来，登记照样能做。\n"
           "这面旗是「我宣布这段内存适合 GDR」，不是「没有这面旗就不能登记」。",
           16, MUTED)
        add_notes(s, """
四步对着 vmm_alloc_and_register() 讲，大约 801 行。
64K 对齐是 GDR_PAGE。REGISTER_VIDMEM 要求地址和长度 64K 对齐，所以 gran 至少 64K。
101 是 CUDA 的 invalid value / not supported 一类，日志会打：cuMemCreate RDMA-capable failed (101), retry without flag。听众在 node2 日志里能看到这句，不是失败结局。
""")

    def s_capture():
        s = add()
        header(s, "06  换分配", "打开「抓 hMemory」窗口：ioctl 和 VMM 在这里会合", len(slides), TOTAL[0])
        code_card(s, 0.48, 1.18, 12.35, 3.85, [
            "capturing_memory = 1;          // 打开窗口",
            "captured_hmemory = 0;",
            "cuMemCreate(...);              // 真 libcuda -> 内核 RM_ALLOC class 0x40",
            "                               // 我们的 ioctl() 偷看成功工单，记下 hNew",
            "hMemory = captured_hmemory;    // 关上窗口，拿走房间号",
            "cuMemMap + cuMemSetAccess;",
            "register_vidmem_locked(va, size, hMemory);",
        ], "vmm_alloc_and_register()  里的会合点")
        bullets(s, 0.55, 5.20, 12.2, 1.7, [
            "没有这个窗口，ioctl 钩子分不清「这一段要给网卡用的显存」和 CUDA 自己的其它小对象。",
            "cuMemCreate / cuMemMap 也被导出，GDRCopy 如果自己走 VMM，同样会登记。",
        ], 16)
        add_notes(s, """
用双手比划：左手是 VMM 替换，右手是 ioctl 记账，在 cuMemCreate 这一拍击掌。
如果听众问「能不能从 VA 反查 hMemory」——理论上有别的 RM 查询，我们选了最稳的：在创建的当口把工单抄下来。
""")

    def s_tpp():
        s = add()
        header(s, "07  补登记本", "自己建一本 0x503c：Tesla 的 libcuda 会做、GeForce 不做", len(slides), TOTAL[0])
        code_card(s, 0.48, 1.18, 7.4, 3.55, [
            "rm_alloc(hClient, hSubdevice,",
            "         0x503c,   /* NV50_THIRD_PARTY_P2P */",
            "         flags = BAR1);",
            "",
            "rm_ctrl(hTPP, 0x503c0102,   /* REGISTER_VA_SPACE */",
            "        { hVASpace });",
            "",
            "rm_ctrl(hTPP, 0x503c0104,   /* REGISTER_VIDMEM */",
            "        { hMemory, va, size });",
        ], "ensure_tpp_locked + register_vidmem_locked")
        card_title(s, 8.10, 1.18, 4.7, 3.55, "对应酒店比喻",
                   "0x503c = 新做一本访客名单\n"
                   "0102 = 声明这本名单管哪个楼层（VA space）\n"
                   "0104 = 把这间房写进名单\n\n"
                   "每个 CUDA client 只建一本。",
                   AMBER, 16, 14)
        tb(s, 0.48, 4.95, 12.3, 1.85,
           "rm_alloc / rm_ctrl 都走 call_ioctl，不再进我们导出的 ioctl()。\n"
           "flags = BAR1：跟 aikitoria 打开的那条物理通道一致。\n"
           "重复登记返回 0x56（INSERT_DUPLICATE_NAME）当成成功——幂等，不怕跑两次。",
           16, SOFT)
        add_notes(s, """
这是「补」字的正文。Tesla 机器上用 strace/GDB 能看到 libcuda 自己发这三下；4090 的 libcuda 日志里永远没有。
parent 必须是 hSubdevice，不是 hDevice。写错父亲，0x503c 会建失败。
BAR1 flag 和我们已有的 P2P 工作对齐：第三方从 BAR1 看 GPU 内存。
0x56 不要当错误讲。
""")

    def s_register_detail():
        s = add()
        header(s, "07  补登记本", "REGISTER_VIDMEM 工单上到底写了什么", len(slides), TOTAL[0])
        box(s, 0.48, 1.22, 6.15, 5.55, BG_CARD, LINE)
        tb(s, 0.72, 1.40, 5.7, 0.4, "字段", 18, AMBER, True)
        rows = [
            ("hMemory", "刚才窗口里抓到的 class 0x40"),
            ("address", "CUDA VA，向下对齐到 64K"),
            ("size", "长度向上补齐到 64K"),
            ("offset", "我们填 0"),
        ]
        y = 1.95
        for a, b in rows:
            tb(s, 0.72, y, 2.1, 0.85, a, 16, WHITE, True, mono=True)
            tb(s, 2.85, y, 3.5, 0.85, b, 15, SOFT)
            y += 0.95
        box(s, 6.85, 1.22, 5.95, 5.55, BG_CODE, LINE)
        tb(s, 7.05, 1.40, 5.55, 0.4, "成功日志", 16, GREEN, True)
        tb(s, 7.05, 1.95, 5.55, 4.5,
           [
               ("REGISTER_VIDMEM", 15, AMBER, True, True),
               ("  hMemory=0x5c00009d", 15, SOFT, False, True),
               ("  va=0x304200000", 15, SOFT, False, True),
               ("  size=0x10000", 15, SOFT, False, True),
               ("  status=0x0", 18, GREEN, True, True),
               ("", 12, MUTED, False, False),
               ("status=0 才算写进本子。", 15, SOFT, False, False),
               ("非 0 且不是 0x56，pin 仍会 -22。", 15, MUTED, False, False),
           ])
        add_notes(s, """
把对齐讲清楚：va & ~(64K-1)，size 把错位的前缀也包进去。否则 RM 直接拒。
hMemory 来自运行时，不要让听众以为示例里的 0x5c00009d 可以写进代码。
""")

    def s_get_pages():
        s = add()
        header(s, "07  补登记本", "为什么这时 gdr_pin_buffer 不再返回 -22", len(slides), TOTAL[0])
        steps = [
            ("应用", "GDRCopy 对刚才那个 VA 调用 gdr_pin_buffer()"),
            ("内核", "走到 nvidia_p2p_get_pages(va, len)"),
            ("RM", "按 VA 去 0x503c 登记本里查"),
            ("命中", "返回 GPU 物理页 / BAR1 映射信息"),
            ("之后", "GDRCopy 才能做 BAR 映射，URMA peermem 才能发"),
        ]
        y = 1.25
        for i, (a, b) in enumerate(steps):
            box(s, 0.48, y, 12.35, 0.95, BG_CARD, LINE, 0.04)
            box(s, 0.68, y + 0.22, 1.7, 0.52, AMBER if i < 4 else GREEN, radius=0.12)
            tb(s, 0.68, y + 0.22, 1.7, 0.52, a, 14, BG, True, align="center",
               anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 2.6, y, 10.0, 0.95, b, 18, WHITE, anchor=MSO_ANCHOR.MIDDLE)
            y += 1.05
        add_notes(s, """
闭环。回到第一张酒店表：保安终于在名单上看到这个门牌。
强调我们没有替换 gdr_pin_buffer，也没有替换 nvidia_p2p_get_pages。库存二进制，手续补齐就能用。
URMA write_lat_gpu --gpu-mode=peermem 依赖的就是这次 pin 成功。
""")

    def s_other_doors():
        s = add()
        header(s, "08  另外两扇门", "属性 116 撒谎，和 SYNC_MEMOPS 的 801", len(slides), TOTAL[0])
        card_title(s, 0.48, 1.22, 6.05, 5.55, "门 A：属性 116 / 110",
                   "Harry 在 libcuda 里把 0x20 拨开。\n"
                   "我们挟持 cuDeviceGetAttribute：\n"
                   "先调真函数，若问的是 116 或 110，把 *pi 改成 1。\n\n"
                   "效果类似，不改文件、不依赖偏移。\n"
                   "驱动升级后不用重算 0x42d69b。\n\n"
                   "只做这一扇，pin 仍会 0x57。",
                   BLUE, 16, 15)
        card_title(s, 6.75, 1.22, 6.05, 5.55, "门 B：SYNC_MEMOPS = 801",
                   "GDRCopy 在 pin 前会 cuPointerSetAttribute(SYNC_MEMOPS)。\n"
                   "4090 上我们对 VMM 指针做这件事，GeForce 返回 801。\n"
                   "GDRCopy 会当成致命错误，即使登记已经成功。\n\n"
                   "钩子：若这个 VA 是我们自己换出来的，把 801 当成 SUCCESS 还回去。\n\n"
                   "日志：owned VMM, return SUCCESS。",
                   AMBER, 16, 15)
        add_notes(s, """
两扇门都是「应用层的礼貌检查」，不是登记本本身。
116：GDRCopy check_gdr_support()。110：VMM 版本的同一句话。
801：CUDA_ERROR_NOT_SUPPORTED。我们只对自己拥有的 VMM VA 撒谎；别人的指针原样返回。
现场可以说：有一次 REGISTER_VIDMEM status=0 了，测试却死在 copylat 的 gpu_mem_alloc，就是这扇门。
""")

    def s_logs():
        s = add()
        header(s, "09  现场", "node2 上成功时，stderr 必须长这样", len(slides), TOTAL[0])
        code_card(s, 0.48, 1.18, 12.35, 4.55, [
            "[gdr-geforce] ready GPU 0 client=... device=... subdev=... va=...",
            "[gdr-geforce] cuMemCreate RDMA-capable failed (101), retry without flag",
            "[gdr-geforce] captured hMemory=0x5c00009d class=0x40",
            "[gdr-geforce] allocated NV50_THIRD_PARTY_P2P hTPP=...",
            "[gdr-geforce] REGISTER_VA_SPACE ... status=0x0",
            "[gdr-geforce] REGISTER_VIDMEM  ... status=0x0",
            "[gdr-geforce] cuMemAlloc -> VMM va=0x304200000",
            "",
            "# 然后 gdr_pin_buffer 不得再是 -22",
        ], "按出现顺序读，就是今天讲的整条链路")
        tb(s, 0.48, 5.90, 12.3, 1.05,
           "缺 ready GPU：ioctl 没套上，或还没 cuInit。\n"
           "有 VMM、无 hMemory：窗口没打开。有登记、pin 仍 -22：pin 的 VA 不在登记范围内。",
           15, MUTED)
        add_notes(s, """
这页当现场对照清单。建议讲座里真的打开一份旧日志念一遍。
101 那行要主动指：不是失败。
va 和 hMemory 每次不同，对数字本身脱敏讲。
""")

    def s_compare():
        s = add()
        header(s, "09  现场", "对照：Harry 的 5090  vs  我们的 4090", len(slides), TOTAL[0])
        headers_row = ("", "Harry · RTX 5090", "我们 · RTX 4090D")
        rows = [
            ("目标", "NCCL + IB dma-buf", "GDRCopy / peermem / get_pages"),
            ("内核", "试过改 dma-buf 能力，不是闸门", "只保留 aikitoria BAR1 P2P"),
            ("用户态手段", "拷一份 libcuda，改一字节", "LD_PRELOAD .so，不改文件"),
            ("拨开的东西", "能力 bit / dma-buf 导出", "attr 116 + 建 NV503C + 登记"),
            ("作者/我们的话", "「这不会让旧 peermem 工作」", "「我们走的就是这条旧路」"),
            ("怎么加载", "LD_LIBRARY_PATH 指向拷贝", "isolate/run.sh，禁止全局 preload"),
        ]
        # header strip
        box(s, 0.40, 1.18, 12.5, 0.55, RGBColor(0x1E, 0x2A, 0x40), None, 0.04)
        tb(s, 0.55, 1.18, 2.4, 0.55, headers_row[0], 14, MUTED, True, anchor=MSO_ANCHOR.MIDDLE)
        tb(s, 3.0, 1.18, 4.7, 0.55, headers_row[1], 15, BLUE, True, anchor=MSO_ANCHOR.MIDDLE)
        tb(s, 7.8, 1.18, 4.9, 0.55, headers_row[2], 15, AMBER, True, anchor=MSO_ANCHOR.MIDDLE)
        y = 1.80
        for i, (a, b, c) in enumerate(rows):
            bg = BG_CARD if i % 2 == 0 else BG_CARD2
            box(s, 0.40, y, 12.5, 0.82, bg, None, 0.03)
            tb(s, 0.55, y, 2.4, 0.82, a, 14, WHITE, True, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 3.0, y, 4.7, 0.82, b, 14, SOFT, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 7.8, y, 4.9, 0.82, c, 14, SOFT, anchor=MSO_ANCHOR.MIDDLE)
            y += 0.82
        add_notes(s, """
对照表是给会后带走的。不要在这里开新战场。
特别指「作者/我们的话」那一行：两篇文章/两份工作并不矛盾，路径不同。
""")

    def s_run():
        s = add()
        header(s, "10  怎么跑", "隔离目录 + run.sh，不要全局 export", len(slides), TOTAL[0])
        code_card(s, 0.48, 1.18, 12.35, 2.85, [
            "./isolate_umdk_gpu.sh build $HOME/umdk_gpu_isolated",
            "$HOME/umdk_gpu_isolated/run.sh \\",
            "  $HOME/umdk_gpu_isolated/bin/urma_perftest write_lat_gpu \\",
            "  --gpu-mode=peermem -d udmac0d1e2 ...",
            "# 临时关掉钩子： UMDK_GDR_HOOK=0  .../run.sh ...",
        ], "只给这一次 exec 设 LD_PRELOAD")
        items = [
            ("为什么不用全局 preload", "sudo 会丢掉 LD_PRELOAD；登录环境里每个进程都会加载 .so。"),
            ("run.sh 会做什么", "make hook、拷进隔离目录、只给测试进程带上 .so，并拦 dlsym。"),
            ("已经是 root 就不要套 sudo", "套了等于把钩子摘掉，然后奇怪地复现「没 ready GPU」。"),
        ]
        y = 4.20
        for t, b in items:
            box(s, 0.48, y, 12.35, 0.85, BG_CARD, LINE, 0.04)
            tb(s, 0.7, y, 3.6, 0.85, t, 15, AMBER, True, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 4.4, y, 8.2, 0.85, b, 15, SOFT, anchor=MSO_ANCHOR.MIDDLE)
            y += 0.90
        add_notes(s, """
操作页。讲座里打开 isolate_umdk_gpu.sh 顶部注释念三句即可。
强调：cat / 打开脚本本身不会加载 so。只有 run.sh exec 出去的那个进程带钩子。
""")

    def s_dont():
        s = add()
        header(s, "10  怎么跑", "不要做的事", len(slides), TOTAL[0])
        rows = [
            ("不要改 4090 的 NVIDIA 内核来「开 GDR」",
             "aikitoria 的 BAR1 P2P 保持原样。GDR 一半在用户态 .so。"),
            ("不要把 Harry 的一字节当成完整方案",
             "那是 dma-buf / NCCL。我们要的是登记本 / get_pages。"),
            ("不要全局 export LD_PRELOAD，不要写 /etc/ld.so.preload",
             "只用 isolate 生成的 run.sh。"),
            ("不要 sudo 掉环境变量之后说钩子坏了",
             "sudo env LD_PRELOAD=绝对路径 ... 或已经是 root 就直接跑。"),
            ("不要写死 GDB 里看到的 handle / 文件偏移",
             "写死 RM 协议。房间号每次运行从 ioctl 抓。"),
        ]
        y = 1.20
        for t, b in rows:
            box(s, 0.48, y, 12.35, 1.08, BG_CARD, CORAL, 0.04)
            tb(s, 0.7, y + 0.10, 12.0, 0.40, t, 16, CORAL, True)
            tb(s, 0.7, y + 0.52, 12.0, 0.45, b, 15, SOFT)
            y += 1.14
        add_notes(s, """
用红色卡片收束纪律。这五条也是 code review 清单。
现场可以点名：谁要是把 hook 写进 /etc/ld.so.preload，第二天登录图形界面都可能怪。sudo 丢掉 LD_PRELOAD 是我们自己复现过的「钩子突然没了」。
写死 handle / 文件偏移，驱动一升级就全废；写死协议 class/cmd 才能跟着 575 活下去。
Harry 的一字节请再指回对照表：不是否定他，是路径不同。
""")

    def s_trouble():
        s = add()
        header(s, "10  怎么跑", "排错表：日志缺哪一句，就回到哪一页", len(slides), TOTAL[0])
        rows = [
            ("从没出现 ready GPU", "ioctl 没套上，或还没 cuInit", "挟持 ioctl / dlsym / run.sh"),
            ("cannot alloc 0x503c yet", "家谱还没收齐就去建本", "track_alloc / 家谱"),
            ("REGISTER_VIDMEM no hMemory", "还在走普通 cuMemAlloc", "VMM 四步 / 抓窗口"),
            ("REGISTER_VIDMEM status != 0", "handle 错，或 VA/size 没 64K 对齐", "登记本页"),
            ("登记成功，pin 仍 -22", "pin 的 VA 不在登记范围内", "对照日志里的 va/size"),
            ("死在 SYNC_MEMOPS / 801", "GDRCopy 把 801 当致命", "另外两扇门"),
        ]
        box(s, 0.40, 1.18, 12.5, 0.50, RGBColor(0x1E, 0x2A, 0x40))
        tb(s, 0.55, 1.18, 4.0, 0.50, "看到的", 13, MUTED, True, anchor=MSO_ANCHOR.MIDDLE)
        tb(s, 4.6, 1.18, 4.3, 0.50, "意思", 13, MUTED, True, anchor=MSO_ANCHOR.MIDDLE)
        tb(s, 9.0, 1.18, 3.7, 0.50, "回到", 13, MUTED, True, anchor=MSO_ANCHOR.MIDDLE)
        y = 1.72
        for i, (a, b, c) in enumerate(rows):
            bg = BG_CARD if i % 2 == 0 else BG_CARD2
            box(s, 0.40, y, 12.5, 0.82, bg)
            tb(s, 0.55, y, 4.0, 0.82, a, 13, WHITE, True, mono=True, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 4.6, y, 4.3, 0.82, b, 14, SOFT, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 9.0, y, 3.7, 0.82, c, 14, AMBER, anchor=MSO_ANCHOR.MIDDLE)
            y += 0.82
        add_notes(s, """
讲座最后 3 分钟用这张表做 Q&A 索引。几乎所有现场问题都能指回前面某一页。
""")

    def s_after():
        s = add()
        header(s, "11  边界", "GDR 打通之后（今天只说到这里）", len(slides), TOTAL[0])
        box(s, 0.48, 1.25, 12.35, 1.70, BG_CARD, GREEN)
        tb(s, 0.75, 1.45, 11.8, 1.35,
           "pin 成功，只说明：网卡侧能按 CUDA VA 找到 GPU 页。\n"
           "把这些页再交给我们自己的 UDMA / 进程页表，是另一场讲座。",
           18, WHITE)
        card_title(s, 0.48, 3.20, 6.05, 3.55, "今天覆盖",
                   "•  用户态闸门\n"
                   "•  ioctl 挟持\n"
                   "•  登记本\n"
                   "•  GDRCopy pin 不再 -22",
                   GREEN, 18, 16)
        card_title(s, 6.75, 3.20, 6.05, 3.55, "今天不覆盖",
                   "•  UMMU / 隔离 SVA\n"
                   "•  MATT / 4K SG\n"
                   "•  进程页表插 PFN\n"
                   "那些是 pin 成功之后的驱动故事。",
                   MUTED, 18, 16)
        add_notes(s, """
主动收住。有人举手问 UMMU 就说：那是下一场，今天的验收标准只有 gdr_pin_buffer 不再 -22，以及 peermem 路径能进到 pin 之后。
不要现场打开 udma 的 PFN 补丁，会把听众带走。
""")

    def s_takeaways():
        s = add()
        header(s, "12  带走", "三句话就够写在笔记本上", len(slides), TOTAL[0])
        cards = [
            ("1", "闸门在用户态",
             "4090 硬件能被 DMA 碰到。GeForce 的 libcuda 不肯宣布、也不肯写登记本。不要再为此改 NVIDIA 内核。"),
            ("2", "Harry 打开的是另一扇门",
             "一字节让 dma-buf / NCCL 活。我们要的是 legacy get_pages。思想借他的，手续借 remote-lab 的。"),
            ("3", "挟持 ioctl 是在补手续",
             "偷看 RM_ALLOC 抄房间号，VMM 逼出 class 0x40，自己建 0x503c 把 VA 写进本。库存 GDRCopy 就能 pin。"),
        ]
        for i, (n, t, b) in enumerate(cards):
            box(s, 0.48, 1.22 + i * 1.85, 12.35, 1.72, BG_CARD, LINE, 0.06)
            box(s, 0.70, 1.48 + i * 1.85, 0.70, 0.70, AMBER, radius=0.16)
            tb(s, 0.70, 1.48 + i * 1.85, 0.70, 0.70, n, 22, BG, True, align="center",
               anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 1.65, 1.40 + i * 1.85, 10.9, 0.45, t, 20, WHITE, True)
            tb(s, 1.65, 1.88 + i * 1.85, 10.9, 0.85, b, 16, SOFT)
        add_notes(s, """
结束前把三句话再念一遍，停顿让人抄：
1) 闸门在用户态，别再改 4090 的 NVIDIA 内核；
2) Harry 打开的是 dma-buf，我们走的是登记本；
3) 挟持 ioctl 是补手续，不是重写 CUDA。
然后翻到代码位置页，告诉想读源码的人从文件头注释开始。
""")

    def s_code_where():
        s = add()
        header(s, "12  带走", "代码在哪，建议怎么读", len(slides), TOTAL[0])
        rows = [
            ("文件头注释", "gdr_geforce_hook.c:1", "先读绑架原理，再读功能列表"),
            ("ioctl / 认人", ":429  :540", "形状过滤 → 真调用 → 拆包"),
            ("track_alloc", ":333", "按 class 记账"),
            ("VMM 替换", ":801  :928", "cuMemAlloc_v2 的身体"),
            ("建本 + 登记", ":631  :671", "0x503c / 0102 / 0104"),
            ("116 与 801", ":910  :949", "两扇小门"),
            ("dlsym", ":1104", "套住 urma_perftest"),
            ("怎么跑", "isolate_umdk_gpu.sh / README.md", "run.sh，不要全局 preload"),
        ]
        y = 1.18
        for a, b, c in rows:
            box(s, 0.48, y, 12.35, 0.68, BG_CARD, LINE, 0.03)
            tb(s, 0.65, y, 2.6, 0.68, a, 14, AMBER, True, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 3.3, y, 4.3, 0.68, b, 14, WHITE, True, mono=True, anchor=MSO_ANCHOR.MIDDLE)
            tb(s, 7.7, y, 4.9, 0.68, c, 14, SOFT, anchor=MSO_ANCHOR.MIDDLE)
            y += 0.70
        add_notes(s, """
给想课后再读代码的人。行号是当前树的近似值，以函数名为准。
建议阅读顺序就是表上的顺序：先文件头，再 ioctl，再 VMM，再建本。
""")

    def s_qa():
        s = add()
        set_slide_bg(s)
        rect(s, 0, 0, 0.16, 7.5, AMBER)
        tb(s, 0.7, 1.8, 12, 0.4, "提问", 16, AMBER, True)
        tb(s, 0.7, 2.3, 12, 1.2, "Q & A", 54, WHITE, True)
        tb(s, 0.7, 3.7, 12, 1.6,
           "想问内核细节？先把它翻译成：\n登记本 / 房间号 / 工单 / 插队，再问。",
           22, SOFT)
        tb(s, 0.7, 5.6, 12, 1.0,
           "源码  userspace/gdr-geforce-hook/\n"
           "文章  harrychen.xyz/2026/05/20/enable-gpudirect-rdma-on-rtx-5090/",
           16, MUTED)
        footer(s, len(slides), TOTAL[0], "提问")
        add_notes(s, """
预留 8–10 分钟。高频问题：
Q: 这算不算攻击 / rootkit？A: 只对自愿 LD_PRELOAD 的测试进程生效，而且只拆 NVIDIA RM 的建对象工单。不要全局安装。
Q: 能不能 upstream 进 NVIDIA？A: GeForce 不开放这条是政策。我们是实验室手续补齐。
Q: 5090 能不能用我们的 so？A: 如果目标是 GDRCopy/peermem，可以试同一套手续；如果目标是 NCCL dma-buf，Harry 的一字节更对症。
Q: 改完 libcuda 一字节再加我们的 so？A: 没必要叠两套。我们的 attr hook 已经覆盖 116。
""")

    # title page uses its own footer-ish, others use TOTAL
    fns = [
        s_title, s_goals, s_contract, s_gdr_picture, s_hotel, s_machine, s_error,
        s_harry_start, s_harry_bit, s_harry_result, s_not_one_byte, s_failures,
        s_borrow, s_bird, s_preload, s_dlsym, s_ioctl_what, s_ioctl_code, s_who,
        s_family, s_track, s_why_vmm, s_vmm_steps, s_capture, s_tpp,
        s_register_detail, s_get_pages, s_other_doors, s_logs, s_compare, s_run,
        s_dont, s_trouble, s_after, s_takeaways, s_code_where, s_qa,
    ]
    TOTAL[0] = len(fns)
    for fn in fns:
        fn()

    # Stamp correct totals on slides that used header() before TOTAL was known.
    # header() already received TOTAL[0] because we set it first.
    prs.save(OUT)
    print(f"wrote {OUT}  ({TOTAL[0]} slides)")


TOTAL = [0]


def main():
    # Fix page numbers: header is called with len(slides) after add(), so page is correct
    # once TOTAL[0] is set before invoking builders. build() does that.
    build()


if __name__ == "__main__":
    main()
